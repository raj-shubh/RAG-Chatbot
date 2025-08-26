from typing import List
from pptx import Presentation
from pptx.util import Pt

from synthesis import Deck, Slide


def _add_title_slide(prs: Presentation, title: str):
    layout = prs.slide_layouts[0]  # Title slide
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    if slide.placeholders and len(slide.placeholders) > 1:
        slide.placeholders[1].text = ""


def _add_bullets_slide(prs: Presentation, title: str, bullets: List[str]):
    layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    body.clear()
    if not bullets:
        p = body.add_paragraph()
        p.text = ""
        return
    # First bullet goes into first paragraph
    body.text = bullets[0]
    # Subsequent bullets as new paragraphs
    for b in bullets[1:]:
        p = body.add_paragraph()
        p.text = b


def build_presentation(deck: Deck, output_path: str) -> str:
    prs = Presentation()
    if not deck.slides:
        _add_title_slide(prs, deck.topic)
    else:
        # Ensure first slide is a title slide
        first = deck.slides[0]
        _add_title_slide(prs, first.title or deck.topic)
        # Remaining slides as bullet slides
        for slide in deck.slides[1:]:
            _add_bullets_slide(prs, slide.title, slide.bullets)
    prs.save(output_path)
    return output_path